"""Generate defense PPT from outline using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === DESIGN TOKENS ===
BLUE = RGBColor(0x1A, 0x3C, 0x6E)
GOLD = RGBColor(0xC8, 0xA0, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)
BLUE_LIGHT = RGBColor(0xE8, 0xEC, 0xF0)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=14, color=DARK, spacing=1.3, font_name='Microsoft YaHei'):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=DARK, bullet='•', font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox

def add_title_bar(slide):
    """Add blue header bar at top."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)
    # Thin gold line below
    add_rect(slide, Inches(0), Inches(0.08), W, Inches(0.02), GOLD)

def add_slide_number(slide, num):
    add_text_box(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
                 str(num), font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_page_title(slide, title, subtitle=None):
    add_title_bar(slide)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 title, font_size=30, color=BLUE, bold=True)
    # Gold underline
    add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.03), GOLD)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=GRAY)

# ============================================================
# SLIDE 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(3.2), BLUE)
add_rect(slide, Inches(0), Inches(3.2), W, Inches(0.06), GOLD)

add_text_box(slide, Inches(1.2), Inches(0.8), Inches(10.5), Inches(1.0),
             '政务数字门户平台 POC', font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(0.6),
             '项目答辩', font_size=28, color=GOLD, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(10.5), Inches(0.5),
             '基于鲲鹏ARM + FusionCompute + Docker 容器化的国产化解决方案', font_size=16, color=RGBColor(0xCC, 0xCC, 0xDD))

add_text_box(slide, Inches(1.2), Inches(4.0), Inches(5), Inches(0.4),
             '小组：实训第4组', font_size=16, color=DARK)
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(5), Inches(0.4),
             '成员：王振光（PM+架构） 胡翰斌 刘永涛 王浩乐', font_size=14, color=GRAY)
add_text_box(slide, Inches(1.2), Inches(5.0), Inches(5), Inches(0.4),
             '日期：2026年7月3日', font_size=14, color=GRAY)

# ============================================================
# SLIDE 2: Project Background
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '客户的痛点', '项目背景')
add_slide_number(slide, 2)

items = [
    ('📄 技术文档分散', '各部门文档格式不统一，检索困难，知识无法沉淀'),
    ('🔒 信创合规压力', '现有系统 x86 架构，不满足国产化要求，政策风险高'),
    ('🚀 交付效率低', '新部门上线需要数周，手动部署环节多，容易出错'),
    ('🛡️ 运维能力弱', '没有监控、没有备份、没有审计，出问题才发现'),
]
for idx, (title, desc) in enumerate(items):
    y = Inches(1.8 + idx * 1.2)
    add_rect(slide, Inches(0.8), y, Inches(0.06), Inches(0.9), GOLD)
    add_text_box(slide, Inches(1.1), y, Inches(10), Inches(0.4), title, font_size=20, color=BLUE, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(10), Inches(0.5), desc, font_size=14, color=GRAY)

# ============================================================
# SLIDE 3: Project Goals
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '我们要做什么', '项目目标')
add_slide_number(slide, 3)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
             '基于 鲲鹏ARM + FusionCompute + Docker 容器技术，搭建一套可复制的政务数字门户 POC 平台',
             font_size=20, color=DARK, bold=True)

goals = [
    '100% 国产化技术栈 — 从芯片到应用全链路自主可控',
    '容器化一键部署 — 单条命令启动，≤30 分钟从零到上线',
    '完整的运维体系 — 监控 + 备份 + 恢复 + 审计 + 快照',
    '可复用模板 — 支撑 10 个部门快速上线新站点',
]
for i, g in enumerate(goals):
    y = Inches(2.8 + i * 0.9)
    add_rect(slide, Inches(1.2), y + Inches(0.05), Inches(0.35), Inches(0.35), GOLD)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(0.35), Inches(0.35),
                 str(i+1), font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y, Inches(10), Inches(0.6), g, font_size=16, color=DARK)

# ============================================================
# SLIDE 4: Team
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '项目团队', '实训第4组')
add_slide_number(slide, 4)

team = [
    ('项目经理+架构', '王振光', '总体负责、系统设计、文档答辩'),
    ('基础设施', '胡翰斌', 'FusionCompute VM、Docker 环境搭建'),
    ('应用开发', '刘永涛', 'Halo 博客、Compose 编排、前后端联调'),
    ('运维', '王浩乐', 'Nginx、监控、备份、一键部署脚本'),
]
for i, (role, name, task) in enumerate(team):
    y = Inches(2.0 + i * 1.3)
    add_rect(slide, Inches(1.5), y, Inches(10), Inches(1.0), LIGHT_BG if i % 2 == 0 else WHITE)
    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(2.5), Inches(0.4), role, font_size=14, color=GOLD, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.5), Inches(2.5), Inches(0.4), name, font_size=20, color=BLUE, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.25), Inches(6.5), Inches(0.5), task, font_size=14, color=GRAY)

# ============================================================
# SLIDE 5: Architecture (core)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '技术架构全景图', '⭐ 核心')
add_slide_number(slide, 5)

# Architecture diagram as text layers
layers = [
    ('用户浏览器 (HTTPS)', '#e8ecf0'),
    ('Nginx — 反向代理 + SSL 终止', '#d0d8e0'),
    ('Halo 博客     MySQL 8.0     Prometheus + Grafana', '#c8d0d8'),
    ('Docker CE 24.x — 容器引擎', '#b8c0c8'),
    ('openEuler 22.03 LTS SP2 (ARM64 aarch64)', '#a8b0b8'),
    ('鲲鹏 TaiShan 200 + FusionCompute 8.x', '#98a0a8'),
]
for i, (label, bg_hex) in enumerate(layers):
    y = Inches(1.6 + i * 0.9)
    rect = add_rect(slide, Inches(2.0), y, Inches(9), Inches(0.7), RGBColor(
        int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16)
    ))
    add_text_box(slide, Inches(2.3), y + Inches(0.15), Inches(8.5), Inches(0.4),
                 label, font_size=14, color=DARK if i < 4 else WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow between layers
    if i < 5:
        add_text_box(slide, Inches(6.2), y + Inches(0.7), Inches(0.8), Inches(0.2),
                     '▼', font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.0), Inches(7.0), Inches(9), Inches(0.3),
             '自下而上：每一层都是国产', font_size=12, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Full-stack Localization
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '100% 国产化技术栈', '🔥 亮点')
add_slide_number(slide, 6)

stack = [
    ('芯片', '鲲鹏 920 (ARM64)', '华为自研'),
    ('虚拟化', 'FusionCompute 8.x', '华为'),
    ('操作系统', 'openEuler 22.03 LTS', '华为开源'),
    ('容器', 'Docker CE 24.x', 'ARM64 官方支持'),
    ('博客', 'Halo v2.x (30K+ ⭐)', '国产开源'),
    ('数据库', 'MySQL 8.0', 'ARM64 官方支持'),
    ('代理', 'Nginx 1.24', 'ARM64 官方支持'),
    ('监控', 'Prometheus + Grafana', 'CNCF 标准'),
]
for i, (layer, tech, source) in enumerate(stack):
    y = Inches(1.6 + i * 0.65)
    add_text_box(slide, Inches(0.8), y, Inches(1.8), Inches(0.4), layer, font_size=13, color=GRAY, bold=True)
    add_text_box(slide, Inches(2.8), y, Inches(4.5), Inches(0.4), tech, font_size=14, color=DARK)
    add_text_box(slide, Inches(7.5), y, Inches(2.5), Inches(0.4), source, font_size=12, color=GOLD)
    add_text_box(slide, Inches(10.3), y, Inches(1.5), Inches(0.4), '✅ 国产', font_size=12, color=GRAY)
    # subtle line
    add_rect(slide, Inches(0.8), y + Inches(0.5), Inches(11.5), Inches(0.005), LIGHT_BG)

# ============================================================
# SLIDE 7: Network Topology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '网络设计', '双网隔离')
add_slide_number(slide, 7)

add_rect(slide, Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8), LIGHT_BG)
add_text_box(slide, Inches(1.3), Inches(1.9), Inches(4.5), Inches(0.4),
             'VLAN 10 — 管理网', font_size=16, color=BLUE, bold=True)
add_text_box(slide, Inches(1.3), Inches(2.3), Inches(4.5), Inches(0.3),
             '192.168.10.0/24', font_size=13, color=GRAY)
add_text_box(slide, Inches(1.3), Inches(2.8), Inches(4.5), Inches(0.3),
             'VM-01 (管理节点)', font_size=15, color=DARK, bold=True)
add_text_box(slide, Inches(1.3), Inches(3.2), Inches(4.5), Inches(0.3),
             '2vCPU / 4GB / 40GB', font_size=12, color=GRAY)
add_text_box(slide, Inches(1.3), Inches(3.6), Inches(4.5), Inches(0.3),
             'Registry 镜像仓库', font_size=13, color=DARK)

add_rect(slide, Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.8), LIGHT_BG)
add_text_box(slide, Inches(7.3), Inches(1.9), Inches(4.5), Inches(0.4),
             'VLAN 20 — 业务网', font_size=16, color=BLUE, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.3), Inches(4.5), Inches(0.3),
             '192.168.20.0/24', font_size=13, color=GRAY)
add_text_box(slide, Inches(7.3), Inches(2.8), Inches(4.5), Inches(0.3),
             'VM-02 (业务节点)', font_size=15, color=DARK, bold=True)
add_text_box(slide, Inches(7.3), Inches(3.2), Inches(4.5), Inches(0.3),
             '2vCPU / 4GB / 40GB + 100GB 数据盘', font_size=12, color=GRAY)
add_bullet_list(slide, Inches(7.3), Inches(3.8), Inches(4.5), Inches(2.5),
                ['Halo 博客', 'MySQL 8.0', 'Nginx (HTTPS)', 'Prometheus + Grafana', 'node_exporter + cadvisor'],
                font_size=13, color=DARK)

add_text_box(slide, Inches(5.8), Inches(3.2), Inches(1.5), Inches(0.5),
             '← 安全组\n   隔离 →', font_size=11, color=GOLD, alignment=PP_ALIGN.CENTER, bold=True)

# ============================================================
# SLIDE 8: Docker Compose
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '容器化部署', 'Docker Compose 编排 7 个服务')
add_slide_number(slide, 8)

services = ['halo (博客) :8090', 'mysql (数据库) :3306', 'nginx (代理) :80/:443',
            'prometheus (采集) :9090', 'grafana (大盘) :3000',
            'node_exporter (主机) :9100', 'cadvisor (容器) :8080']
for i, s in enumerate(services):
    col = i % 2
    row = i // 2
    x = Inches(1.2 + col * 5.8)
    y = Inches(1.8 + row * 1.3)
    add_rect(slide, x, y, Inches(5.2), Inches(1.0), LIGHT_BG)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.5), Inches(0.3),
                 f'容器 {i+1}', font_size=10, color=GOLD, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(4.5), Inches(0.3),
                 s, font_size=15, color=DARK, bold=True)

add_rect(slide, Inches(3.0), Inches(6.0), Inches(7), Inches(0.7), BLUE)
add_text_box(slide, Inches(3.2), Inches(6.1), Inches(6.5), Inches(0.5),
             '$ docker compose up -d          ← 一条命令启动全部', font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: Live Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '现场演示', '🎬 核心环节 — 6 个步骤')
add_slide_number(slide, 9)

demos = [
    'SSH 登录 VM-02 → uname -m 确认 aarch64 ARM 架构',
    'docker ps 检查 7 个容器状态 — 全部 UP',
    '浏览器访问 https://192.168.20.10 → 展示博客首页',
    '发布一篇测试文章 → 草稿 → 审核 → 前台可见',
    'docker rm -f mysql → 重建容器 → 数据不丢失验证',
    'Grafana 监控大盘 — CPU/内存/容器实时图表',
]
for i, d in enumerate(demos):
    y = Inches(1.7 + i * 0.85)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), LIGHT_BG if i % 2 == 0 else WHITE)
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.4),
                 str(i+1), font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.6), y + Inches(0.12), Inches(10.5), Inches(0.4),
                 d, font_size=14, color=DARK)

# ============================================================
# SLIDE 10: Review Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '内容安全 — 审核工作流', '政务合规要求')
add_slide_number(slide, 10)

flow = [
    ('作者发布\n(草稿)', BLUE),
    ('提交审核\n(待审核)', GOLD),
    ('审核员通过 ✅\n→ 前台可见 (已发布)', RGBColor(0x2E, 0x7D, 0x32)),
    ('审核员驳回 ❌\n→ 退回作者 (草稿)', RGBColor(0xC6, 0x28, 0x28)),
]
for i, (label, color) in enumerate(flow):
    x = Inches(1.0 + i * 3.0)
    add_rect(slide, x, Inches(2.5), Inches(2.5), Inches(2.0), color)
    add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.1), Inches(1.2),
                 label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, x + Inches(2.5), Inches(3.2), Inches(0.5), Inches(0.5),
                     '→', font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: One-click Deploy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '30 分钟从零到上线', '🔥 亮点 — 一键部署能力')
add_slide_number(slide, 11)

deploy_steps = [
    ('1/6', '检查 Docker 环境', '~1 分钟'),
    ('2/6', '拉取 ARM64 镜像', '~5 分钟'),
    ('3/6', '生成 SSL 证书', '~10 秒'),
    ('4/6', '创建数据目录', '~5 秒'),
    ('5/6', '启动 7 个容器', '~3 分钟'),
    ('6/6', '等待服务就绪检查', '~2 分钟'),
]
for i, (step, desc, time) in enumerate(deploy_steps):
    y = Inches(1.7 + i * 0.8)
    add_text_box(slide, Inches(0.8), y, Inches(0.8), Inches(0.4), step, font_size=12, color=GOLD, bold=True)
    add_text_box(slide, Inches(1.8), y, Inches(6), Inches(0.4), desc, font_size=15, color=DARK)
    # time bar
    bar_w = [1, 5, 0.2, 0.1, 3, 2][i] * 0.8
    add_rect(slide, Inches(8.0), y + Inches(0.1), Inches(bar_w), Inches(0.25), BLUE)
    add_text_box(slide, Inches(8.0 + bar_w + 0.2), y, Inches(2), Inches(0.4), time, font_size=12, color=GRAY)

add_rect(slide, Inches(3.0), Inches(6.5), Inches(7), Inches(0.6), BLUE)
add_text_box(slide, Inches(3.2), Inches(6.55), Inches(6.5), Inches(0.5),
             '总计 ~28 分钟 ✅ (目标 ≤30min)', font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Ops System
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '不只是能跑，更是可维护', '运维体系')
add_slide_number(slide, 12)

ops_cards = [
    ('📊 监控大盘', 'Grafana + Prometheus\nCPU / 内存 / 容器 / 网络\n实时图表 + 告警规则'),
    ('💾 自动备份', 'backup.sh + crontab\n每日 02:00 执行\n保留最近 7 天 + 每周归档'),
    ('🔄 一键恢复', '备份 → 新环境恢复\n完整验证流程\nRTO ≤ 30min / RPO ≤ 24h'),
    ('📋 日志审计', 'SSH 登录日志\n命令操作记录\n保留 90 天可追溯'),
    ('📸 VM 快照', 'FusionCompute 快照\n4 个关键节点\n可随时快速回滚'),
]
for i, (title, desc) in enumerate(ops_cards):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(1.6 + (i // 3) * 2.8)
    add_rect(slide, x, y, Inches(3.8), Inches(2.4), LIGHT_BG)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4),
                 title, font_size=16, color=BLUE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(1.5),
                 desc, font_size=12, color=GRAY)

# ============================================================
# SLIDE 13: Performance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '性能验证', '实测数据')
add_slide_number(slide, 13)

metrics = [
    ('页面响应 (95分位)', '≤ 3s', '2.1s', '✅'),
    ('并发用户', '50', '50', '✅'),
    ('错误率', '0%', '0%', '✅'),
    ('TPS', '≥ 20', '28', '✅'),
    ('恢复时间 (RTO)', '≤ 30min', '28min', '✅'),
    ('数据恢复 (RPO)', '≤ 24h', '24h', '✅'),
]
add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5), BLUE)
for j, h in enumerate(['指标', '目标', '实测', '状态']):
    add_text_box(slide, Inches(1.0 + j * 3.2), Inches(1.55), Inches(2.5), Inches(0.4),
                 h, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
for i, (metric, target, actual, status) in enumerate(metrics):
    y = Inches(2.1 + i * 0.75)
    if i % 2 == 0:
        add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), LIGHT_BG)
    for j, val in enumerate([metric, target, actual, status]):
        c = GRAY if j == 0 else (GOLD if j == 3 else DARK)
        add_text_box(slide, Inches(1.0 + j * 3.2), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     val, font_size=14, color=c, bold=(j > 0), alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 14: Security
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '安全设计', '7 层防护')
add_slide_number(slide, 14)

security = [
    'HTTPS 全站加密 — TLS 1.2 / 1.3',
    '80 → 443 强制跳转 — HSTS 策略',
    '管理网 / 业务网隔离 — 安全组拦截',
    '数据库仅业务网可访问 — 无公网暴露',
    'SSH 登录 + 命令操作审计 — 保留 90 天',
    '文章发布审核工作流 — 敏感信息把关',
    '容器最小权限运行 — 非 root 用户',
]
for i, s in enumerate(security):
    y = Inches(1.6 + i * 0.8)
    add_rect(slide, Inches(1.0), y, Inches(0.4), Inches(0.4), GOLD if i < 4 else BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(0.4), Inches(0.3),
                 '🔒', font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(10), Inches(0.4),
                 s, font_size=15, color=DARK)

# ============================================================
# SLIDE 15: Highlights
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '三大核心竞争力', '项目亮点总结')
add_slide_number(slide, 15)

hl = [
    ('1', '全栈国产化', '鲲鹏ARM + FusionCompute + openEuler + Halo\n→ 完全自主可控，满足信创要求'),
    ('2', '容器化可复制', 'docker compose up -d + deploy.sh\n→ 30分钟内部署到任意鲲鹏环境'),
    ('3', '运维完备性', '监控 + 备份 + 恢复 + 审计 + 快照\n→ 不是"能跑就行"，是"有人维护"'),
]
for i, (num, title, desc) in enumerate(hl):
    y = Inches(1.8 + i * 1.8)
    add_rect(slide, Inches(1.0), y, Inches(11), Inches(1.5), LIGHT_BG)
    add_rect(slide, Inches(1.0), y, Inches(0.8), Inches(1.5), GOLD)
    add_text_box(slide, Inches(1.0), y + Inches(0.4), Inches(0.8), Inches(0.6),
                 num, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.2), y + Inches(0.2), Inches(3), Inches(0.4),
                 title, font_size=22, color=BLUE, bold=True)
    add_text_box(slide, Inches(2.2), y + Inches(0.7), Inches(9), Inches(0.6),
                 desc, font_size=13, color=GRAY)

# ============================================================
# SLIDE 16: Deliverables
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '交付物清单', '项目文档交付')
add_slide_number(slide, 16)

docs = [
    ('📄', '项目需求分析', '基于 BRD → SRS 转化'),
    ('📄', '系统架构设计', '架构图 + 技术栈 + 设计决策'),
    ('📄', '部署实施手册', '从零到上线的每一步操作'),
    ('📄', '运维操作手册', '日常运维 + 故障处理 SOP'),
    ('📊', '华为项目管理模板', '策划 / 进度 / 会议纪要 ×4'),
    ('💻', '源代码 + 脚本', 'docker-compose + deploy + backup'),
    ('🎬', '答辩 PPT', '本演示文档'),
]
for i, (icon, name, desc) in enumerate(docs):
    col = i % 2
    row = i // 2
    x = Inches(1.0 + col * 6.0)
    y = Inches(1.6 + row * 1.5)
    add_rect(slide, x, y, Inches(5.5), Inches(1.2), LIGHT_BG if row % 2 == 0 else WHITE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.4),
                 icon, font_size=20, color=GOLD)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.15), Inches(4.4), Inches(0.4),
                 name, font_size=16, color=BLUE, bold=True)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.6), Inches(4.4), Inches(0.4),
                 desc, font_size=12, color=GRAY)

# ============================================================
# SLIDE 17: Summary & Learnings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '收获与踩坑', '项目总结')
add_slide_number(slide, 17)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.4),
             '💡 技术层面', font_size=18, color=BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
                ['ARM 架构下容器化部署完整流程', 'Docker Compose 多服务编排', 'Prometheus + Grafana 监控体系', '国产化全栈技术验证'], font_size=13, color=DARK)

add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
             '🔧 工程层面', font_size=18, color=BLUE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.5),
                ['BRD 需求 → SRS 技术转化流程', '策划→执行→监控→交付项目管理', '文档驱动的知识转移', '团队协作与分工效率'], font_size=13, color=DARK)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.4),
             '⚠️ 踩过的坑', font_size=18, color=GOLD, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.5),
                ['Docker Hub 国内拉取慢 → 配置阿里云镜像加速 + 离线镜像包备份',
                 'ARM64 镜像兼容性 → 提前在鲲鹏环境验证所有镜像，避免 x86 假设',
                 'SSL 自签名证书浏览器警告 → POC 阶段可接受，生产使用 Let\'s Encrypt'],
                font_size=13, color=GRAY)

# ============================================================
# SLIDE 18: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(3.0), BLUE)
add_rect(slide, Inches(0), Inches(3.0), W, Inches(0.06), GOLD)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
             '感谢聆听，请老师指正！', font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.6),
             'Q & A', font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             '政务数字门户平台 POC', font_size=22, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
             '基于鲲鹏ARM + 容器化的国产化解决方案', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.6),
             '实训第4组  ·  王振光 胡翰斌 刘永涛 王浩乐  ·  2026年7月3日', font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), '答辩PPT.pptx')
prs.save(output_path)
print(f'PPT saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
