"""Generate defense PPT for Yushe Blog + Digital Moss project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === DESIGN ===
DARK   = RGBColor(0x1E, 0x1B, 0x18)  # warm dark
ACCENT = RGBColor(0x8B, 0x6F, 0x4E)  # warm brown-gold
LIGHT  = RGBColor(0xFB, 0xF8, 0xF3)  # warm cream
GRAY   = RGBColor(0xA8, 0x98, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x6B, 0x8C, 0x42)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def bg(slide, c): slide.background.fill.solid(); slide.background.fill.fore_color.rgb = c
def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h); s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s
def tb(slide, l, t, w, h, txt, sz=18, c=DARK, b=False, a=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = fn; p.alignment = a; return bx
def mtext(slide, l, t, w, h, lines, sz=14, c=DARK, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = fn; p.space_after = Pt(4)
    return bx
def bullets(slide, l, t, w, h, items, sz=13, c=DARK, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = fn; p.space_after = Pt(6)
    return bx
def pagenum(slide, n):
    tb(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3), str(n), sz=10, c=GRAY, a=PP_ALIGN.RIGHT)
def titlebar(slide, title, sub=None):
    rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)
    tb(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55), title, sz=28, c=DARK, b=True)
    rect(slide, Inches(0.8), Inches(0.95), Inches(1.0), Inches(0.025), ACCENT)
    if sub: tb(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.35), sub, sz=13, c=GRAY)

# ===== S1: COVER =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
rect(s, Inches(0), Inches(0), W, Inches(3.2), DARK)
rect(s, Inches(0), Inches(3.2), W, Inches(0.05), ACCENT)
tb(s, Inches(1), Inches(0.7), Inches(11), Inches(0.9), '虞舍 · 数字苔藓', sz=44, c=LIGHT, b=True)
tb(s, Inches(1), Inches(1.6), Inches(11), Inches(0.5), '个人独立博客 + 活的网页花园', sz=24, c=ACCENT, b=True)
tb(s, Inches(1), Inches(2.3), Inches(11), Inches(0.4), '单文件全栈架构 · 自建部署 · AI 驱动的互动体验', sz=15, c=GRAY)
tb(s, Inches(1), Inches(4.2), Inches(6), Inches(0.4), '实训第4组   王振光', sz=16, c=DARK, b=True)
tb(s, Inches(1), Inches(4.7), Inches(6), Inches(0.3), '2026年7月3日', sz=13, c=GRAY)

# ===== S2: PROJECT OVERVIEW =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '项目全景', '我们做了什么')
pagenum(s, 2)
cards = [('📝 虞舍博客', '111 篇原创文章\n4 个分类 / 标签云 / RSS\n搜索 / 归档 / 热榜 / 海报', True),
         ('🌿 数字苔藓', '活的网页花园\n6 种植物 / 四季 / 昼夜\n浇水互动 / 萤火虫 / 蜗牛', False),
         ('🔧 后端 API', 'Python 单文件 HTTP Server\n博客 CRUD + 评论 + RSS\n花园状态 + 浇水端点', True),
         ('🚀 自建部署', '校内服务器 10.42.78.75\nNginx 反向代理\nPodman 容器化', False)]
for i, (t, d, hl) in enumerate(cards):
    x = Inches(0.5 + (i % 2) * 6.3); y = Inches(1.6 + (i // 2) * 2.6)
    rect(s, x, y, Inches(6.0), Inches(2.3), LIGHT if not hl else WHITE)
    tb(s, x + Inches(0.3), y + Inches(0.2), Inches(5.4), Inches(0.4), t, sz=18, c=DARK, b=True)
    mtext(s, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1.3), d.split('\n'), sz=13, c=GRAY)

# ===== S3: TECH STACK =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '技术栈', '无框架 · 零依赖 · 纯手写')
pagenum(s, 3)
stack = [('前端', 'HTML5 + CSS3 + Canvas 2D\n单文件架构，零 JS 框架\n响应式 + 暗黑模式'),
         ('后端', 'Python http.server\nJSON 文件数据库\nRESTful API 设计'),
         ('部署', '校内 CentOS 服务器\nNginx 反向代理\nPodman 容器 + SSH 运维'),
         ('AI 集成', '智谱 GLM-4V 视觉识别\nQwen-VL 场景分析\nDepth Anything V2 深度图')]
for i, (t, d) in enumerate(stack):
    x = Inches(0.6 + i * 3.1); y = Inches(1.7)
    rect(s, x, y, Inches(2.9), Inches(4.8), WHITE)
    tb(s, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.4), t, sz=16, c=ACCENT, b=True, a=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.5), y + Inches(0.65), Inches(1.9), Inches(0.02), ACCENT)
    mtext(s, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(3.5), d.split('\n'), sz=12, c=GRAY)

# ===== S4: BLOG ARCHITECTURE =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '虞舍博客架构', 'yushe-blog.html + wall_api.py')
pagenum(s, 4)
layers = [('浏览器 (HTTPS)', DARK),
          ('Nginx :80 → 反向代理 /api/ → :8089', ACCENT),
          ('wall_api.py (Python HTTP Server :8089)', RGBColor(0x5A, 0x4A, 0x3A)),
          ('JSON 文件存储 (blog_articles.json)', GRAY)]
for i, (label, color) in enumerate(layers):
    y = Inches(1.6 + i * 1.1)
    rr = rect(s, Inches(2.5), y, Inches(8.3), Inches(0.8), color)
    tb(s, Inches(2.8), y + Inches(0.15), Inches(7.8), Inches(0.5), label, sz=15, c=WHITE if color != LIGHT else DARK, b=True, a=PP_ALIGN.CENTER)
features = ['🔍 全文搜索 / 标签筛选', '📅 时间线归档', '🔥 热门文章 TOP 10', '⏳ 时间胶囊（定时发布）',
            '🗺️ Canvas 鸟瞰图', '🎨 文章海报生成', '💬 评论系统', '📡 RSS 订阅', '🌙 暗黑模式']
for i, f in enumerate(features):
    col = i % 3; row = i // 3
    tb(s, Inches(1.0 + col * 3.9), Inches(5.5 + row * 0.5), Inches(3.6), Inches(0.4), f, sz=11, c=GRAY)

# ===== S5: DIGITAL MOSS =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), GREEN)
tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55), '数字苔藓 — 活的网页花园', sz=28, c=LIGHT, b=True)
tb(s, Inches(0.8), Inches(0.95), Inches(11), Inches(0.35), 'moss.html  |  642 行 Canvas 2D  |  零依赖  |  程序化生成', sz=13, c=GRAY)
pagenum(s, 5)

gfx = [('🌓 昼夜循环', '真实时间驱动\n天空色温渐变\n萤火虫夜间更亮'),
       ('🌱 6 种植物', '蕨 / 花 / 蘑菇 / 草\n多肉 / 藤蔓\n程序化绘制，每株不同'),
       ('💧 浇水互动', '点击/触摸浇水\n粒子雨滴反馈\n健康值 + 生长动画'),
       ('🍂 四季系统', '春华 / 夏茂 / 秋实 / 冬雪\n植物生长速度变化\n天空 + 地面换季'),
       ('🦋 活物生态', '蜗牛爬行 + 黏液轨迹\n蝴蝶随机飞动\n萤火虫夜晚闪烁'),
       ('📋 木牌语录', '10 条轮播语录\n双击晃落树叶\n25 秒自动切换')]
for i, (t, d) in enumerate(gfx):
    x = Inches(0.4 + (i % 3) * 4.2); y = Inches(1.6 + (i // 3) * 2.7)
    rect(s, x, y, Inches(3.9), Inches(2.4), RGBColor(0x2A, 0x25, 0x20))
    tb(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.35), t, sz=15, c=GREEN, b=True)
    mtext(s, x + Inches(0.2), y + Inches(0.65), Inches(3.5), Inches(1.5), d.split('\n'), sz=12, c=RGBColor(0xA8, 0x98, 0x80))

# ===== S6: BLOG FEATURES =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '博客功能矩阵', '111 篇文章 · 全功能单文件博客')
pagenum(s, 6)
cols_data = [
    ['内容管理', ['Markdown 文章系统', '4 个分类 (技术/生活/随笔/分享)', '标签云 + 标签筛选', '文章置顶', '封面图支持', '定时发布 / 时间胶囊']],
    ['阅读体验', ['全文搜索 (/)', '时间线归档视图', '热门文章 TOP 10', 'Canvas 鸟瞰图 (可视化)', '文章海报生成下载', '暗黑模式切换']],
    ['互动系统', ['文章评论', 'RSS 订阅', '友链页面', '随机阅读 🎲', 'Ctrl+K 管理面板', '代码块一键复制']],
]
for ci, (col_title, items) in enumerate(cols_data):
    x = Inches(0.6 + ci * 4.1)
    tb(s, x, Inches(1.5), Inches(3.8), Inches(0.4), col_title, sz=16, c=ACCENT, b=True)
    rect(s, x, Inches(2.0), Inches(3.6), Inches(0.02), ACCENT)
    for i, item in enumerate(items):
        tb(s, x + Inches(0.1), Inches(2.2 + i * 0.55), Inches(3.5), Inches(0.45), f'• {item}', sz=12, c=GRAY)

# ===== S7: BACKEND API =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '后端 API 设计', 'wall_api.py — 单文件 Python HTTP Server (564 行)')
pagenum(s, 7)
ep = [('GET /blog/articles', '返回全部文章 JSON'),
      ('POST /blog/articles', '发布文章 (管理员)'),
      ('PUT /blog/articles/:id', '更新文章'),
      ('DELETE /blog/articles/:id', '删除文章'),
      ('POST /blog/articles/:id/comments', '发表评论'),
      ('POST /blog/articles/:id/view', '阅读计数'),
      ('GET /blog/rss', 'RSS Feed'),
      ('GET /garden/state', '花园状态 (植物+季节)'),
      ('POST /garden/visit', '访问播种 (新植物)'),
      ('POST /garden/water', '浇水 (含限流)')]
for i, (method, desc) in enumerate(ep):
    col = i // 5; row = i % 5
    x = Inches(0.6 + col * 6.3); y = Inches(1.5 + row * 1.05)
    rect(s, x, y, Inches(5.9), Inches(0.85), WHITE if row % 2 == 0 else LIGHT)
    tb(s, x + Inches(0.15), y + Inches(0.08), Inches(3.8), Inches(0.3), method, sz=11, c=ACCENT, b=True)
    tb(s, x + Inches(0.15), y + Inches(0.45), Inches(5.5), Inches(0.3), desc, sz=12, c=GRAY)

# ===== S8: GARDEN DATA MODEL =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '花园数据模型', 'garden.json — 植物生命周期')
pagenum(s, 8)
tb(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4), '数据字段', sz=16, c=ACCENT, b=True)
fields = ['id: 唯一标识 (时间戳+随机)', 'type: 植物类型 (fern/flower/mushroom/grass/succulent/vine/glowing)',
          'x, y: 归一化坐标 (0~1)', 'size: 尺寸 (0.3~1.0)', 'health: 健康值 (0~100)',
          'stage: 种子→嫩芽→成长→开花→枯萎', 'createdAt / lastWateredAt: 时间戳',
          'wateredBy[]: 浇水者列表']
for i, f in enumerate(fields):
    col = i // 4; row = i % 4
    tb(s, Inches(0.8 + col * 6.3), Inches(2.1 + row * 0.55), Inches(5.8), Inches(0.4), f'• {f}', sz=12, c=GRAY)
tb(s, Inches(0.8), Inches(4.6), Inches(11), Inches(0.4), '核心机制', sz=16, c=ACCENT, b=True)
mech = ['访问播种: 每次访问自动生成 1 株新植物 (sessionStorage 去重)', '浇水照料: 点击植物 → health+20 → 生长动画 → 粒子雨滴',
        '自然衰减: -0.5 health/小时 → 枯萎 (灰度+透明)', '季节影响: 春季花多、秋季蘑菇爆发、冬季减速+飘雪',
        '稀有事件: 连续7天浇水 → 发光植物；秋季15%蘑菇爆发', '容量控制: 上限 200 株，优先移除枯萎的']
for i, m in enumerate(mech):
    tb(s, Inches(0.8), Inches(5.1 + i * 0.35), Inches(11.5), Inches(0.3), f'• {m}', sz=11, c=GRAY)

# ===== S9: DEPLOYMENT =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '部署架构', '校内服务器 10.42.78.75')
pagenum(s, 9)
layers2 = [('用户浏览器', DARK),
           ('Nginx :80 (Podman 容器)', ACCENT),
           ('/ → yushe-blog.html    /moss → wall_api.py    /garden/ → wall_api.py    /moss-assets/ → wall_api.py', RGBColor(0x5A, 0x4A, 0x3A)),
           ('wall_api.py (:8089) — HTTP API Server', RGBColor(0x7A, 0x6A, 0x5A)),
           ('JSON 文件存储 (/opt/blog/data/)', GRAY)]
for i, (label, color) in enumerate(layers2):
    y = Inches(1.5 + i * 0.95)
    rr = rect(s, Inches(2.0), y, Inches(9.3), Inches(0.7), color)
    tb(s, Inches(2.3), y + Inches(0.12), Inches(8.8), Inches(0.45), label, sz=13, c=WHITE if color != LIGHT else DARK, b=True, a=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4), '部署命令:  scp → ssh kill → nohup python3 wall_api.py &  |  nginx -s reload', sz=12, c=GRAY)

# ===== S10: INNOVATION =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
rect(s, Inches(0), Inches(0), W, Inches(0.05), GREEN)
tb(s, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55), '创新点', sz=28, c=LIGHT, b=True)
pagenum(s, 10)
innovations = [
    ('1', '活的页面 — 数字苔藓', '网页不再是静态的——它会生长、枯萎、呼吸。每次访问改变花园生态，时间本身就是交互方式。这是对"网页"定义的重新想象。'),
    ('2', '单文件全栈', '博客 635 行 HTML，后端 564 行 Python，花园 642 行 Canvas。没有框架、没有构建工具、没有 npm install。极致简单，极致可控。'),
    ('3', '程序化生成美学', '6 种植物纯 Canvas 绘制，18 层渲染管线（天空→远山→土丘→雾→生物→植物→暗角），四季昼夜自然过渡，萤火虫/蜗牛/蝴蝶让花园有生命感。'),
    ('4', 'AI 驱动的博客生态', '智谱 GLM-4V 免费视觉识别，Qwen-VL 场景分析，Depth Anything V2 深度图——AI 不是噱头，是实际可用的功能模块。'),
]
for i, (num, title, desc) in enumerate(innovations):
    y = Inches(1.5 + i * 1.45)
    rect(s, Inches(0.6), y, Inches(12.1), Inches(1.25), RGBColor(0x2A, 0x25, 0x20))
    rect(s, Inches(0.6), y, Inches(0.7), Inches(1.25), GREEN)
    tb(s, Inches(0.6), y + Inches(0.35), Inches(0.7), Inches(0.5), num, sz=28, c=DARK, b=True, a=PP_ALIGN.CENTER)
    tb(s, Inches(1.6), y + Inches(0.1), Inches(3.5), Inches(0.35), title, sz=16, c=LIGHT, b=True)
    tb(s, Inches(1.6), y + Inches(0.55), Inches(10.8), Inches(0.55), desc, sz=11, c=GRAY)

# ===== S11: COMPARISON =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '为什么不用现成方案', '我们的选择 vs 常规做法')
pagenum(s, 11)
comp = [
    ('维度', '常规做法', '我们的做法', '优势'),
    ('博客系统', 'WordPress / Halo', '手写单文件 HTML', '零维护，极快加载'),
    ('后端', 'Spring Boot / Express', 'Python http.server', '564 行，部署无依赖'),
    ('数据库', 'MySQL / PostgreSQL', 'JSON 文件', '零配置，备份=复制文件'),
    ('花园', '不存在', 'Canvas 2D 程序化', '原创交互体验'),
    ('前端框架', 'React / Vue', '纯 HTML+CSS+JS', '0 依赖，1 个文件'),
    ('部署', 'Docker + K8s', 'scp + nohup', '一行命令上线'),
]
xs = [Inches(0.6), Inches(2.8), Inches(6.0), Inches(9.5)]
ws = [Inches(2.2), Inches(3.2), Inches(3.5), Inches(3.5)]
for i, row in enumerate(comp):
    y = Inches(1.5 + i * 0.85)
    for j, val in enumerate(row):
        c = WHITE if i == 0 else (ACCENT if j == 3 else (DARK if j == 2 else GRAY))
        tb(s, xs[j], y + Inches(0.15), ws[j], Inches(0.5), val, sz=11 if i > 0 else 12, c=c, b=(i==0))

# ===== S12: LEARNINGS =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
titlebar(s, '收获与反思', '从零到上线的完整闭环')
pagenum(s, 12)
learn = [
    ('💡 技术成长', ['Python HTTP Server 从入门到生产', 'Canvas 2D 程序化渲染完整实践', 'Nginx 反向代理 + Podman 容器', '无框架单文件架构的设计哲学']),
    ('🔧 工程能力', ['从 BRD → 设计 → 开发 → 部署 → 运维', 'Git 工作流 (feature branch + PR + push)', 'SSH 远程部署 + 进程管理', 'JSON 数据存储的读写优化']),
    ('⚠️ 踩过的坑', ['jsdelivr CDN 国内阻塞 → vendor/ 本地副本', 'Firebase Auth 需走代理 → 端口直连', 'GFW 阻断 GitHub → Clash 代理配置', 'Canvas 中文渲染模糊 → 缩放适配']),
    ('🎯 产品思维', ['不是"能用就行"，是"有人愿意用"', '博客 111 篇文章 + 花园互动 = 完整生态', '每次迭代都有明确交付物', '文档驱动：设计→计划→实现→部署']),
]
for i, (title, items) in enumerate(learn):
    x = Inches(0.4 + (i % 2) * 6.4); y = Inches(1.5 + (i // 2) * 2.8)
    tb(s, x + Inches(0.2), y, Inches(5), Inches(0.35), title, sz=15, c=ACCENT, b=True)
    for j, item in enumerate(items):
        tb(s, x + Inches(0.2), y + Inches(0.5 + j * 0.42), Inches(5.8), Inches(0.35), f'• {item}', sz=11, c=GRAY)

# ===== S13: THANK YOU =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, LIGHT)
rect(s, Inches(0), Inches(0), W, Inches(2.8), DARK)
rect(s, Inches(0), Inches(2.8), W, Inches(0.05), ACCENT)
tb(s, Inches(1), Inches(0.5), Inches(11), Inches(0.8), '感谢聆听，请老师指正', sz=36, c=LIGHT, b=True, a=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(1.5), Inches(11), Inches(0.5), 'Q & A', sz=26, c=ACCENT, b=True, a=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(4.0), Inches(11), Inches(0.6), '虞舍 · 数字苔藓', sz=22, c=DARK, b=True, a=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(4.6), Inches(11), Inches(0.5), '个人独立博客 + 活的网页花园  |  单文件全栈架构  |  自建部署', sz=14, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(5.2), Inches(11), Inches(0.4), '实训第4组  ·  王振光  ·  2026年7月3日', sz=12, c=GRAY, a=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(5.8), Inches(11), Inches(0.4), 'http://10.42.78.75  |  http://10.42.78.75/moss', sz=12, c=ACCENT, a=PP_ALIGN.CENTER)

# ===== SAVE =====
out = os.path.join(os.path.dirname(__file__), '答辩PPT.pptx')
try:
    prs.save(out)
except PermissionError:
    out = os.path.join(os.path.dirname(__file__), '答辩PPT_new.pptx')
    prs.save(out)
print(f'Saved: {out} ({len(prs.slides)} slides)')
